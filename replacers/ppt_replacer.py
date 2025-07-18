from pptx import Presentation
import win32com.client

def replace_in_ppt(file_path, replacements):
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for old, new in replacements:
                    shape.text = shape.text.replace(old, new)
    # 母版
    for master in prs.masters:
        for shape in master.shapes:
            if shape.has_text_frame:
                for old, new in replacements:
                    shape.text = shape.text.replace(old, new)
    prs.save(file_path)

def replace_in_ppt_ppt(file_path, replacements):
    ppt = win32com.client.Dispatch('PowerPoint.Application')
    ppt.Visible = True
    pres = ppt.Presentations.Open(file_path, WithWindow=False)
    for slide in pres.Slides:
        for shape in slide.Shapes:
            if shape.HasTextFrame:
                for old, new in replacements:
                    if shape.TextFrame.HasText:
                        shape.TextFrame.TextRange.Replace(old, new)
    pres.Save()
    pres.Close()
    ppt.Quit() 